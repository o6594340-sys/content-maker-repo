# -*- coding: utf-8 -*-
from flask import Flask, request, jsonify, render_template
import requests
import os
import tempfile

app = Flask(__name__)

# Загружаем .env вручную
def _load_env():
    env_path = os.path.join(os.path.dirname(__file__), ".env")
    if not os.path.exists(env_path):
        return
    with open(env_path, encoding="utf-8") as f:
        for line in f:
            line = line.strip()
            if line and not line.startswith("#") and "=" in line:
                key, _, val = line.partition("=")
                os.environ[key.strip()] = val.strip()

_load_env()

from prompts import CHANNELS, CONTENT_PLAN_SYSTEM, CONTENT_PLAN_CHANNEL_CONTEXT

AI_PROVIDER = os.getenv("AI_PROVIDER", "claude")


@app.route("/")
def index():
    channels = [{"id": k, "name": v["name"], "description": v["description"]} for k, v in CHANNELS.items()]
    provider = AI_PROVIDER
    return render_template("index.html", channels=channels, provider=provider)


@app.route("/generate", methods=["POST"])
def generate():
    data = request.json
    channel_id = data.get("channel")
    topic = data.get("topic", "").strip()
    brief = data.get("brief", "").strip()
    post_type = data.get("post_type", "post")

    if not channel_id or channel_id not in CHANNELS:
        return jsonify({"error": "Канал не выбран"}), 400
    if post_type == "rewrite":
        if not brief:
            return jsonify({"error": "Вставьте свой текст или тезисы в поле ниже"}), 400
    elif not topic:
        return jsonify({"error": "Укажите тему поста"}), 400

    channel = CHANNELS[channel_id]
    system_prompt = channel["system_prompt"]
    user_prompt = _build_user_prompt(topic, brief, post_type)

    try:
        if AI_PROVIDER == "claude":
            result = _call_claude(system_prompt, user_prompt)
        elif AI_PROVIDER == "openai":
            result = _call_openai(system_prompt, user_prompt)
        elif AI_PROVIDER == "gigachat":
            result = _call_gigachat(system_prompt, user_prompt)
        else:
            result = _call_ollama(system_prompt, user_prompt)
        return jsonify({"text": result})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route("/set_provider", methods=["POST"])
def set_provider():
    global AI_PROVIDER
    data = request.json
    provider = data.get("provider", "claude")
    if provider in ("claude", "openai", "gigachat", "ollama"):
        AI_PROVIDER = provider
    return jsonify({"provider": AI_PROVIDER})


@app.route("/content_plan", methods=["POST"])
def content_plan():
    channel_id = request.form.get("channel", "all")
    manual_text = request.form.get("text", "").strip()
    file = request.files.get("file")

    extracted = ""

    if file and file.filename:
        filename = file.filename.lower()
        with tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(filename)[1]) as tmp:
            file.save(tmp.name)
            tmp_path = tmp.name
        try:
            if filename.endswith(".pptx"):
                extracted = _parse_pptx(tmp_path)
            elif filename.endswith(".pdf"):
                extracted = _parse_pdf(tmp_path)
            elif filename.endswith(".docx"):
                extracted = _parse_docx(tmp_path)
            else:
                return jsonify({"error": "Неподдерживаемый формат. Используй PPTX, PDF или DOCX."}), 400
        finally:
            os.unlink(tmp_path)

    if manual_text:
        extracted = (extracted + "\n\n" + manual_text).strip() if extracted else manual_text

    if not extracted:
        return jsonify({"error": "Загрузи файл или вставь текст"}), 400

    channel_ctx = CONTENT_PLAN_CHANNEL_CONTEXT.get(channel_id, CONTENT_PLAN_CHANNEL_CONTEXT["all"])
    user_prompt = f"Канал: {channel_ctx}\n\nМатериал:\n{extracted[:12000]}"

    try:
        if AI_PROVIDER == "claude":
            result = _call_claude(CONTENT_PLAN_SYSTEM, user_prompt)
        elif AI_PROVIDER == "openai":
            result = _call_openai(CONTENT_PLAN_SYSTEM, user_prompt)
        elif AI_PROVIDER == "gigachat":
            result = _call_gigachat(CONTENT_PLAN_SYSTEM, user_prompt)
        else:
            result = _call_ollama(CONTENT_PLAN_SYSTEM, user_prompt)
        return jsonify({"text": result})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


def _parse_pptx(path):
    from pptx import Presentation
    prs = Presentation(path)
    lines = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in para.runs).strip()
                    if text:
                        lines.append(text)
        lines.append("")
    return "\n".join(lines).strip()


def _parse_pdf(path):
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    pages = []
    for page in reader.pages:
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text.strip())
    return "\n\n".join(pages)


def _parse_docx(path):
    import zipfile
    from xml.etree import ElementTree as ET
    ns = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
    with zipfile.ZipFile(path) as z:
        xml = z.read("word/document.xml")
    root = ET.fromstring(xml)
    lines = []
    for p in root.iter(ns + "p"):
        text = "".join(r.text or "" for r in p.iter(ns + "t")).strip()
        lines.append(text)
    return "\n".join(lines).strip()


@app.route("/health")
def health():
    return jsonify({"status": "ok", "provider": AI_PROVIDER})


def _build_user_prompt(topic, brief, post_type):
    if post_type == "rewrite":
        return (
            f"Вот мой черновик / тезисы:\n\n{brief}\n\n"
            "Перепиши это в готовый пост для Telegram в стиле канала. "
            "Сохрани мои мысли и факты, но сделай текст живым, цепляющим и в нужном стиле. "
            "Не добавляй то, чего нет в исходнике."
        )
    type_instructions = {
        "post": "Напиши полноценный пост для Telegram.",
        "ideas": "Предложи 5 идей для постов по этой теме. Для каждой — заголовок и 1-2 строки о чём.",
        "short": "Напиши короткий пост для Telegram (до 80 слов).",
    }
    instruction = type_instructions.get(post_type, type_instructions["post"])
    prompt = f"Тема: {topic}\n\n"
    if brief:
        prompt += f"Дополнительный контекст / тезисы:\n{brief}\n\n"
    prompt += instruction
    return prompt


def _call_claude(system_prompt, user_prompt):
    import anthropic
    client = anthropic.Anthropic(api_key=os.getenv("ANTHROPIC_API_KEY"))
    message = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=1024,
        system=system_prompt,
        messages=[{"role": "user", "content": user_prompt}],
    )
    return message.content[0].text


def _call_openai(system_prompt, user_prompt):
    from openai import OpenAI
    client = OpenAI(api_key=os.getenv("OPENAI_API_KEY"))
    response = client.chat.completions.create(
        model="gpt-4o-mini",
        messages=[
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        max_tokens=1024,
        temperature=0.8,
    )
    return response.choices[0].message.content


def _call_gigachat(system_prompt, user_prompt):
    import base64
    client_id = os.getenv("GIGACHAT_CLIENT_ID", "")
    client_secret = os.getenv("GIGACHAT_CLIENT_SECRET", "")
    credentials = base64.b64encode(f"{client_id}:{client_secret}".encode()).decode()

    # Получаем токен
    token_resp = requests.post(
        "https://ngw.devices.sberbank.ru:9443/api/v2/oauth",
        headers={"Authorization": f"Basic {credentials}", "Content-Type": "application/x-www-form-urlencoded"},
        data={"scope": "GIGACHAT_API_PERS"},
        verify=False,
        timeout=30,
    )
    token_resp.raise_for_status()
    token = token_resp.json()["access_token"]

    # Генерируем текст
    resp = requests.post(
        "https://gigachat.devices.sberbank.ru/api/v1/chat/completions",
        headers={"Authorization": f"Bearer {token}", "Content-Type": "application/json"},
        json={
            "model": "GigaChat",
            "messages": [
                {"role": "system", "content": system_prompt},
                {"role": "user", "content": user_prompt},
            ],
            "max_tokens": 1024,
            "temperature": 0.8,
        },
        verify=False,
        timeout=60,
    )
    resp.raise_for_status()
    return resp.json()["choices"][0]["message"]["content"]


def _call_ollama(system_prompt, user_prompt):
    payload = {
        "model": os.getenv("OLLAMA_MODEL", "llama3.2:3b"),
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        "stream": False,
        "options": {"temperature": 0.8, "top_p": 0.9},
    }
    response = requests.post(
        f"{os.getenv('OLLAMA_URL', 'http://localhost:11434')}/api/chat",
        json=payload,
        timeout=300,
    )
    response.raise_for_status()
    return response.json()["message"]["content"]


if __name__ == "__main__":
    app.run(host="0.0.0.0", port=5000, debug=True)
